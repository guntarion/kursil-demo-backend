# app/services/document_service.py
import os
import re
import logging
from docx import Document
from docx.shared import Inches, Pt
from bson import ObjectId
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from markdown import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
from datetime import datetime
from ..db.operations import main_topic_collection, list_topics_collection, points_discussion_collection

logger = logging.getLogger(__name__)

async def get_all_handouts_by_main_topic_id(main_topic_id: str):
    # Get the main topic document
    main_topic = await main_topic_collection.find_one({"_id": ObjectId(main_topic_id)})
    if not main_topic:
        return None, []

    all_handouts = []

    # Iterate through list_of_topics
    for topic_name in main_topic['list_of_topics']:
        # Find the corresponding document in list_topics_collection
        list_topic = await list_topics_collection.find_one({"topic_name": topic_name})
        if not list_topic:
            continue

        # Get all points_discussion documents for this topic
        points_discussion_cursor = points_discussion_collection.find({
            "topic_name_id": list_topic['_id']
        })
        points_discussion_docs = await points_discussion_cursor.to_list(length=None)

        # Collect non-empty handouts
        for point in points_discussion_docs:
            if point.get('handout'):
                all_handouts.append({
                    'topic_name': topic_name,
                    'point_of_discussion': point['point_of_discussion'],
                    'handout': point['handout']
                })

    return main_topic['main_topic'], all_handouts

def sanitize_filename(filename):
    """Remove or replace characters that are unsafe for filenames."""
    # Replace spaces and other unsafe characters with underscores
    filename = re.sub(r'[^\w\-_\. ]', '_', filename)
    # Replace multiple spaces with a single underscore
    filename = re.sub(r'\s+', '_', filename)
    return filename


async def generate_handout_word_document(main_topic: str, handouts: list):
    doc = Document()
    doc.add_heading('Handouts Document', 0)
    print("main_topic === ", main_topic)
    # Get existing styles
    styles = doc.styles

    # Group handouts by topic_name, maintaining original order
    current_topic = None
    for handout in handouts:
        if handout['topic_name'] != current_topic:
            if current_topic is not None:
                doc.add_page_break()
            current_topic = handout['topic_name']
            doc.add_heading(current_topic, level=1)
        
        doc.add_heading(handout['point_of_discussion'], level=2)
        
        # Convert markdown to HTML
        html = markdown(handout['handout'])
        
        # Parse HTML
        soup = BeautifulSoup(html, 'html.parser')
        
        for element in soup.find_all():
            if element.name == 'h4':
                doc.add_paragraph(element.text, style='Heading 3')
            elif element.name == 'p':
                paragraph = doc.add_paragraph()
                for child in element.children:
                    if child.name == 'strong':
                        paragraph.add_run(child.text).bold = True
                    else:
                        paragraph.add_run(child.text)
            elif element.name == 'ul':
                for li in element.find_all('li'):
                    doc.add_paragraph(li.text, style='List Bullet')

    # Add final page break if there were any handouts
    if handouts:
        doc.add_page_break()

    # Generate filename
    date_str = datetime.now().strftime("%y-%m-%d")
    sanitized_topic = sanitize_filename(main_topic)
    base_filename = f"{sanitized_topic} - {date_str} - handout"

    # Ensure the documents directory exists
    os.makedirs('./documents', exist_ok=True)

    # Check for existing files and append number if necessary
    index = 0
    while True:
        if index == 0:
            filename = f"{base_filename}.docx"
        else:
            filename = f"{base_filename} ({index}).docx"

        document_path = os.path.join('./documents', filename)
        if not os.path.exists(document_path):
            break
        index += 1

    doc.save(document_path)
    return document_path



async def get_all_kursil_by_main_topic_id(main_topic_id: str):
    logger.info(f"Fetching kursil data for main topic ID: {main_topic_id}")
    
    # Get the main topic document
    main_topic = await main_topic_collection.find_one({"_id": ObjectId(main_topic_id)})
    if not main_topic:
        logger.warning(f"Main topic not found for ID: {main_topic_id}")
        return None, []

    logger.info(f"Main topic found: {main_topic.get('main_topic', 'Unnamed Main Topic')}")

    all_kursil_data = []

    # Get all list_topics documents for this main topic
    # Try with both ObjectId and string versions of main_topic_id
    list_topics_cursor = list_topics_collection.find({
        "$or": [
            {"main_topic_id": ObjectId(main_topic_id)},
            {"main_topic_id": main_topic_id}
        ]
    })
    list_topics_docs = await list_topics_cursor.to_list(length=None)
    
    logger.info(f"Found {len(list_topics_docs)} list topics for main topic")

    if len(list_topics_docs) == 0:
        # If no documents found, let's log some additional information
        sample_doc = await list_topics_collection.find_one()
        if sample_doc:
            logger.info(f"Sample document from list_topics_collection: {sample_doc}")
        else:
            logger.warning("list_topics_collection appears to be empty")

        all_main_topic_ids = await list_topics_collection.distinct("main_topic_id")
        logger.info(f"All main_topic_ids in list_topics_collection: {all_main_topic_ids}")

    for list_topic in list_topics_docs:
        logger.info(f"Processing list topic: {list_topic.get('topic_name', 'Unnamed Topic')}")
        logger.info(f"List topic ID: {list_topic['_id']}")
        logger.info(f"List topic main_topic_id: {list_topic.get('main_topic_id')}")
        
        topic_data = {
            'topic_name': list_topic.get('topic_name', 'Unnamed Topic'),
            'objective': list_topic.get('objective', ''),
            'key_concepts': list_topic.get('key_concepts', ''),
            'skills_to_be_mastered': list_topic.get('skills_to_be_mastered', ''),
            'points_of_discussion': list_topic.get('point_of_discussion', []),
            'details': []
        }

        # Get all points_discussion documents for this topic
        points_discussion_cursor = points_discussion_collection.find({
            "topic_name_id": list_topic['_id']
        })
        points_discussion_docs = await points_discussion_cursor.to_list(length=None)
        
        logger.info(f"Found {len(points_discussion_docs)} points of discussion for topic {topic_data['topic_name']}")

        for point in points_discussion_docs:
            logger.info(f"+== Processing point of discussion: {point.get('point_of_discussion', 'Unnamed Point')}")
            logger.info(f"Point of discussion ID: {point['_id']}")
            logger.info(f"Point of discussion: {point.get('point_of_discussion', '')}")
            logger.info(f"learn_objective: {point.get('learn_objective', '')}")
            logger.info(f"assessment: {point.get('assessment', '')}")
            logger.info(f"method: {point.get('method', '')}")
            
            point_data = {
                'point_of_discussion': point.get('point_of_discussion', ''),
                'learn_objective': point.get('learn_objective', ''),
                'assessment': point.get('assessment', ''),
                'method': point.get('method', ''),
                'duration': point.get('duration', '')
            }
            topic_data['details'].append(point_data)

        all_kursil_data.append(topic_data)

    logger.info(f"Total kursil data points collected: {len(all_kursil_data)}")
    return main_topic.get('main_topic', 'Unnamed Main Topic'), all_kursil_data


async def generate_kursil_word_document(main_topic: str, kursil_data: list):
    doc = Document()
    
    # Add or get styles
    styles = doc.styles
    if 'Heading 4' not in styles:
        style = styles.add_style('Heading 4', WD_STYLE_TYPE.PARAGRAPH)
        style.font.size = Pt(12)
        style.font.bold = True
    else:
        style = styles['Heading 4']
        style.font.size = Pt(12)
        style.font.bold = True
    
    # Title
    title = doc.add_heading('Kursil Document', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Main topic
    doc.add_heading(main_topic, level=1)
    doc.add_paragraph()  # Add some space

    for topic_data in kursil_data:
        logger.info(f"Processing topic: {topic_data.get('topic_name', 'Unnamed Topic')}")
        doc.add_heading(topic_data.get('topic_name', 'Unnamed Topic'), level=2)

        # Add topic information from list_topics_collection
        if topic_data.get('objective'):
            doc.add_heading('Objective', level=3)
            doc.add_paragraph(topic_data['objective'])

        if topic_data.get('key_concepts'):
            doc.add_heading('Key Concepts', level=3)
            doc.add_paragraph(topic_data['key_concepts'])

        if topic_data.get('skills_to_be_mastered'):
            doc.add_heading('Skills to be Mastered', level=3)
            doc.add_paragraph(topic_data['skills_to_be_mastered'])

        if topic_data.get('points_of_discussion'):
            doc.add_heading('Points of Discussion', level=3)
            for point in topic_data['points_of_discussion']:
                doc.add_paragraph(point, style='List Bullet')

        # Add details from points_discussion_collection
        for detail in topic_data.get('details', []):
            if detail.get('point_of_discussion'):
                doc.add_heading(detail['point_of_discussion'], level=3)

            for field in ['learn_objective', 'assessment', 'method']:
                content = detail.get(field, '')
                if content:
                    doc.add_heading(field.replace('_', ' ').title(), level=4)

                    # Convert markdown to HTML
                    html = markdown(content)

                    # Parse HTML
                    soup = BeautifulSoup(html, 'html.parser')

                    for element in soup.find_all():
                        if element.name == 'h4':
                            doc.add_paragraph(element.text, style='Heading 4')
                        elif element.name == 'p':
                            paragraph = doc.add_paragraph()
                            for child in element.children:
                                if child.name == 'strong':
                                    paragraph.add_run(child.text).bold = True
                                else:
                                    paragraph.add_run(child.text)
                        elif element.name in ['ul', 'ol']:
                            for li in element.find_all('li', recursive=False):
                                # Remove any existing numbering or bullet points
                                text = re.sub(r'^\d+\.\s*|\•\s*|\-\s*', '', li.text.strip())
                                p = doc.add_paragraph(text, style='List Bullet')
                                # Ensure the paragraph is part of the list
                                p.paragraph_format.left_indent = Inches(0.25)
                                p.paragraph_format.first_line_indent = Inches(-0.25)
                            # Add an empty paragraph to end the list
                            doc.add_paragraph()

            # Duration
            if detail.get('duration'):
                doc.add_heading('Duration', level=4)
                doc.add_paragraph(f"{detail['duration']} minutes")

        doc.add_page_break()

    # Generate filename
    date_str = datetime.now().strftime("%y-%m-%d")
    sanitized_topic = sanitize_filename(main_topic)
    base_filename = f"{sanitized_topic} - {date_str} - kursil"

    # Ensure the documents directory exists
    os.makedirs('./documents', exist_ok=True)

    # Check for existing files and append number if necessary
    index = 0
    while True:
        if index == 0:
            filename = f"{base_filename}.docx"
        else:
            filename = f"{base_filename} ({index}).docx"

        document_path = os.path.join('./documents', filename)
        if not os.path.exists(document_path):
            break
        index += 1

    doc.save(document_path)
    logger.info(f"Document saved at: {document_path}")
    return document_path

async def get_powerpoint_data_by_main_topic_id(main_topic_id: str):
    # Get the main topic document
    main_topic = await main_topic_collection.find_one({"_id": ObjectId(main_topic_id)})
    if not main_topic:
        return None, []

    powerpoint_data = []

    # Iterate through list_of_topics
    for topic_name in main_topic['list_of_topics']:
        # Find the corresponding document in list_topics_collection
        list_topic = await list_topics_collection.find_one({"topic_name": topic_name})
        if not list_topic:
            continue

        # Get all points_discussion documents for this topic
        points_discussion_cursor = points_discussion_collection.find({
            "topic_name_id": list_topic['_id']
        })
        points_discussion_docs = await points_discussion_cursor.to_list(length=None)

        # Collect non-empty elaborations
        for point in points_discussion_docs:
            if point.get('elaboration'):
                powerpoint_data.append({
                    'topic_name': topic_name,
                    'point_of_discussion': point['point_of_discussion'],
                    'elaboration': point['elaboration']
                })

    return main_topic['main_topic'], powerpoint_data



async def generate_powerpoint_document(main_topic: str, powerpoint_data: list):
    # Create a presentation with 16:9 aspect ratio
    # prs = Presentation()
    # prs.slide_width = Inches(16)
    # prs.slide_height = Inches(9)

    # Load the custom template
    template_path = os.path.join('./app/res/', 'template_plnnp.pptx')
    prs = Presentation(template_path)

    # Title slide (assuming it's the first layout in your template)
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Adjust these lines based on your template's placeholder indices
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = main_topic
    if subtitle:
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"

    for data in powerpoint_data:
        # Section header slide for point_of_discussion
        section_slide_layout = prs.slide_layouts[2]  # Assuming layout 2 is "Section Header"
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = data['point_of_discussion']

        # Parse the elaboration content
        content_lines = data['elaboration'].split('\n')
        current_title = ""
        current_content = []

        for line in content_lines:
            if line.startswith('- **') and line.endswith('**'):
                # If we have a previous title and content, create a slide for it
                if current_title and current_content:
                    create_content_slide(prs, current_title, current_content)
                
                # Start a new title and reset content
                current_title = line.strip('- **')
                current_title = current_title.strip(':')
                current_content = []
            elif line.strip():
                body_content = line.strip('- ')
                current_content.append(body_content)

        # Create the last slide if there's remaining content
        if current_title and current_content:
            create_content_slide(prs, current_title, current_content)

    # Generate filename
    date_str = datetime.now().strftime("%y-%m-%d")
    sanitized_topic = sanitize_filename(main_topic)
    base_filename = f"{sanitized_topic} - {date_str} - presentation"

    # Ensure the documents directory exists
    os.makedirs('./documents', exist_ok=True)

    # Check for existing files and append number if necessary
    index = 0
    while True:
        if index == 0:
            filename = f"{base_filename}.pptx"
        else:
            filename = f"{base_filename} ({index}).pptx"

        document_path = os.path.join('./documents', filename)
        if not os.path.exists(document_path):
            break
        index += 1

    prs.save(document_path)
    return document_path

def create_content_slide(prs, title, content):
    content_slide_layout = prs.slide_layouts[1]  # Assuming layout 1 is "Title and Content"
    slide = prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = title

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = ""

    for point in content:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0